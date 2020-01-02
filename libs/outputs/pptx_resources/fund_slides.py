from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import pandas as pd 
import numpy as np 
import os 
import glob 
import datetime

from libs.utils import fund_list_extractor, windows_compatible_file_parse, SP500
from libs.tools import trend_simple_forecast

from .slide_utils import slide_title_header, color_to_RGB

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8


def make_fund_slides(prs, analysis: dict):
    funds = analysis.keys()
    for fund in funds:
        prs = add_fund_content(prs, fund, analysis)

    return prs


def add_fund_content(prs, fund: str, analysis: dict):
    content_dir = f'output/temp/{fund}/'
    if os.path.exists(content_dir):
        # Title slide for a fund
        fund_name = SP500.get(fund, fund)

        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        top = Inches(0.1)
        left = Inches(4)
        width = Inches(5)
        height = Inches(2)
        txtbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txtbox.text_frame

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f'{fund_name}'
        p.font.bold = True
        p.font.size = Pt(60)
        p.font.name = 'Arial'

        p2 = text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = f"Dates Covered: {analysis[fund]['dates_covered']['start']}  :  {analysis[fund]['dates_covered']['end']}"
        p2.font.bold = False
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x74, 0x3c, 0xe6)
        p2.font.name = 'Arial'

        has_beta = False
        if 'beta' in analysis[fund].keys():
            # Insert a table of fund figures
            left_loc = Inches(0.1)
            top_loc = Inches(1.1)
            table_width = Inches(2.4)
            table_height = Inches(1.4)

            vq = analysis[fund].get('metadata', {}).get('volatility', {}).get('VQ')
            has_vq = False
            rows = 3
            if vq is not None:
                has_vq = True
                rows = 6
                stop_loss = analysis[fund].get('metadata', {}).get('volatility', {}).get('stop_loss')
                print(f"vol: {analysis[fund].get('metadata', {}).get('volatility', {}).get('last_max')}")
                high_close = analysis[fund].get('metadata', {}).get('volatility', {}).get('last_max', {}).get('Price', 'n/a')

            table_placeholder = slide.shapes.add_table( rows, 
                                                        2,
                                                        left_loc,
                                                        top_loc,
                                                        table_width,
                                                        table_height)
            table = table_placeholder.table

            table.cell(0,0).text = 'Attribute'
            table.cell(0,1).text = ''
            table.cell(1,0).text = 'Beta'
            table.cell(1,1).text = str(np.round(analysis[fund]['beta'], 5))
            table.cell(2,0).text = 'R-Squared'
            table.cell(2,1).text = str(np.round(analysis[fund]['r_squared'], 5))

            if has_vq:
                table.cell(3,0).text = 'Volatility Quotient'
                table.cell(3,1).text = str(vq)
                table.cell(4,0).text = 'Stop Loss'
                table.cell(4,1).text = str(stop_loss)
                table.cell(5,0).text = 'Last High Close'
                table.cell(5,1).text = str(high_close)

            table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(16)
            table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(16)
            for i in range(1, rows):
                table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(14)
            has_beta = True

        content = content_dir + f"candlestick_{fund}.png"
        if os.path.exists(content):
            if has_beta:
                left = Inches(2.6) #Inches(1.42)
            else:
                left = Inches(1.42)
            top = Inches(1.4)
            height = Inches(6)
            width = Inches(10.5)
            slide.shapes.add_picture(content, left, top, height=height, width=width)

        price_pt = np.round(analysis[fund]['statistics']['current_price'],2)
        price_chg_p = np.round(analysis[fund]['statistics']['current_percent_change'],3)
        price_chg = np.round(analysis[fund]['statistics']['current_change'],2)
        if price_chg > 0.0:
            price_str = f"${price_pt} +{price_chg} (+{price_chg_p}%)"
        else:
            price_str = f"${price_pt} {price_chg} ({price_chg_p}%)"

        # Slide #1 of content
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        indexes = []
        indexes.append(len(prs.slides) - 1)
        slide = slide_title_header(slide, fund, price_details=price_str)

        # Slide #2 of content
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, fund, price_details=price_str)
        indexes.append(len(prs.slides)-1)

        # Slide #3 of content
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, fund, price_details=price_str)
        indexes.append(len(prs.slides)-1)

        # Slide #4 of content
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, fund, price_details=price_str)
        indexes.append(len(prs.slides)-1)

        # Slide #5 of content
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
        slide = slide_title_header(slide, fund, price_details=price_str)
        indexes.append(len(prs.slides)-1)

        content = content_dir + '*.png'
        pics = glob.glob(content)
        fund_analysis = analysis[fund]
        prs = format_plots(prs, indexes, pics, fund_analysis=fund_analysis)

    return prs


def format_plots(prs, slide_indices: list, globs: list, fund_analysis: dict={}):
    parts = windows_compatible_file_parse(globs[0])

    header = parts[0] + '/' + parts[1] + '/' + parts[2] + '/'

    for globber in globs:

        globbed = windows_compatible_file_parse(globber)
        part = globbed[3]

        if 'volume' in part:
            left = Inches(0)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'RSI_standard' in part:
            left = Inches(6.5)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'macd_bar' in part:
            left = Inches(0.0)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'clustering' in part:
            left = Inches(6.5)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[0]].shapes.add_picture(header+part, left, top, height=height, width=width)

        ### Slide #2

        if 'relative_strength' in part:
            left = Inches(0)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'exp_moving_averages' in part:
            left = Inches(6.5)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'swing_trades' in part:
            left = Inches(0.0)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'simple_moving_averages' in part:
            left = Inches(6.5)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[1]].shapes.add_picture(header+part, left, top, height=height, width=width)

        # Slide #3

        if 'obv_standard' in part:
            left = Inches(0)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[2]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'obv_diff' in part:
            left = Inches(6.5)
            top = Inches(1.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[2]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'head_and_shoulders' in part:
            left = Inches(0.0)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[2]].shapes.add_picture(header+part, left, top, height=height, width=width)

        if 'price_gaps' in part:
            left = Inches(6.5)
            top = Inches(4.1)
            height = Inches(3.0)
            width = Inches(6.5)
            prs.slides[slide_indices[2]].shapes.add_picture(header+part, left, top, height=height, width=width)

        # Slide #4

        if 'resist_support' in part:
            left = Inches(0)
            top = Inches(1.55)
            height = Inches(4.7)
            width = Inches(7)
            prs.slides[slide_indices[3]].shapes.add_picture(header+part, left, top, height=height, width=width)

            left = Inches(7)
            top = Inches(0.25)
            height = Inches(4.7)
            width = Inches(4)
            txbox = prs.slides[slide_indices[3]].shapes.add_textbox(left, top, width, height)
            
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Nearest Support & Resistance Levels"
            p.font.size = Pt(18)
            p.font.name = 'Arial'
            p.font.bold = True

            left_loc = Inches(8)
            top_loc = Inches(1)
            table_width = Inches(4)

            num_srs = len(fund_analysis['support_resistance']['major S&R']) + 1
            table_height = Inches(num_srs * 0.33)
            if num_srs * 0.33 > 6.0:
                table_height = Inches(6.0)

            table_placeholder = prs.slides[slide_indices[3]].shapes.add_table(
                                                    num_srs, 
                                                    2,
                                                    left_loc,
                                                    top_loc,
                                                    table_width,
                                                    table_height)
            table = table_placeholder.table

            table.cell(0,0).text = 'Price'
            table.cell(0,1).text = '% Change'

            for i, maj in enumerate(fund_analysis['support_resistance']['major S&R']):
                table.cell(i+1, 0).text = f"${maj['Price']}"
                table.cell(i+1, 1).text = f"{maj['Change']}"
                color = color_to_RGB(maj['Color'])
                table.cell(i+1, 0).text_frame.paragraphs[0].font.color.rgb = color
                table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = color

        # Slide 5

        if 'trendline' in part:
            left = Inches(0)
            top = Inches(1.55)
            height = Inches(4.7)
            width = Inches(7)
            prs.slides[slide_indices[4]].shapes.add_picture(header+part, left, top, height=height, width=width)

            left = Inches(6.5)
            top = Inches(0.25)
            height = Inches(4.7)
            width = Inches(4)
            txbox = prs.slides[slide_indices[4]].shapes.add_textbox(left, top, width, height)
            
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Long, Intermediate, Short, and Near Term Trendlines"
            p.font.size = Pt(18)
            p.font.name = 'Arial'
            p.font.bold = True

            trends = []
            futures = list(range(0,91,15))
            forecasts = []
            for trend in fund_analysis['trendlines']:
                if trend['current']:
                    trends.append(trend)
                    forecast = trend_simple_forecast(trend, future_periods=futures)
                    forecasts.append(forecast)

            num_rows = len(trends) + 2
            num_cols = len(futures)
            table_height = num_rows * 0.33
            table_width = num_cols * 0.8
            if num_rows * 0.33 > 6.0:
                table_height = 6.0
            if num_cols * 0.8 > 6.6:
                table_width = 6.6

            left_start = 6.5
            adj_left = (6.6 - table_width) / 2
            table_height = Inches(table_height)
            table_width = Inches(table_width)
            left_loc = Inches(left_start + adj_left)
            top_loc = Inches(0.75)

            table_placeholder = prs.slides[slide_indices[4]].shapes.add_table(
                                                    num_rows, 
                                                    num_cols,
                                                    left_loc,
                                                    top_loc,
                                                    table_width,
                                                    table_height)
            table = table_placeholder.table

            cell_1 = table.cell(0,0)
            cell_2 = table.cell(0,num_cols-1)
            cell_1.merge(cell_2)
            cell_1.text = f"Future Periods of Active Trendlines"

            for i, fut in enumerate(futures):
                table.cell(1,i).text = str(fut)
                table.cell(1,i).text_frame.paragraphs[0].font.bold = True

            for i, trend in enumerate(trends):
                for j, value in enumerate(forecasts[i]['returns']):
                    table.cell(i+2,j).text = f"${value}"
                    table.cell(i+2,j).text_frame.paragraphs[0].font.size = Pt(12)
                    color = color_to_RGB(trend['color'])
                    table.cell(i+2,j).text_frame.paragraphs[0].font.color.rgb = color


    return prs 