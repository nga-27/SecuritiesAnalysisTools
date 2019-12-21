from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from datetime import datetime
import os 

from libs.utils import SP500

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8


def slide_title_header(slide, fund: str, include_time=True, price_details=''):
    fund = SP500.get(fund, fund)
    
    left = Inches(0) #Inches(3.86)
    top = Inches(0)
    width = height = Inches(0.5)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame

    p = tf.paragraphs[0]
    p.text = fund 
    p.font.size = Pt(36)
    p.font.name = 'Arial'
    p.font.bold = True

    if include_time:
        p = tf.add_paragraph()
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.name = 'Arial'
        p.text = str(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
    
    if price_details != '':
        left = Inches(2) 
        top = Inches(0.1)
        width = height = Inches(0.5)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame

        p = tf.paragraphs[0]
        p.text = fund 
        p.font.size = Pt(24)
        p.font.name = 'Arial'
        p.font.bold = True
        p.text = price_details
        deets = price_details.split(' ')
        if '+' in deets[1]:
            p.font.color.rgb = color_to_RGB('green')
        else:
            p.font.color.rgb = color_to_RGB('red')


    return slide


def subtitle_header(slide, title: str):
    """ Creates subtitle under main slide title """
    top = Inches(0.61)
    left = Inches(0.42)
    width = height = Inches(0.5)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txtbox.text_frame 

    p = text_frame.paragraphs[0]
    p.text = title 
    p.font.bold = False 
    p.font.size = Pt(22)
    p.font.name = 'Times New Roman'

    return slide


def intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, 'Explanation of Analysis', include_time=False)

    if os.path.exists('resources/metric_explanation.txt'):
        filer = open('resources/metric_explanation.txt', 'r')
        content = filer.readlines()
        content2 = []
        for cont in content:
            c = cont.split('\r\n')[0]
            content2.append(c)
        content = content2
        filer.close()

        top = Inches(0.81)
        left = Inches(0.42)
        width = Inches(11)
        height = Inches(6)
        txtbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txtbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = content[0] 
        p.font.size = Pt(12)
        p.font.bold = True
        for i in range(1,len(content)):
            p = text_frame.add_paragraph()
            p.text = content[i]
            if i == 3:
                p.font.size = Pt(12)
                p.font.bold = True
            else:
                p.font.size = Pt(10)
                p.font.bold = False

    else:
        print("WARNING - file 'metric_explanation.txt' not found.")

    return prs


def color_to_RGB(color: str):
    if color == 'black':
        return RGBColor(0x00, 0x00, 0x00)
    elif color == 'blue':
        return RGBColor(0x00, 0x00, 0xFF)
    elif color == 'green':
        return RGBColor(0x00, 0xCC, 0x00)
    elif color == 'purple':
        return RGBColor(0x7F, 0x00, 0xFF)
    elif color == 'yellow':
        return RGBColor(0xee, 0xd4, 0x00)
    elif color == 'orange':
        return RGBColor(0xff, 0x99, 0x33)
    elif color == 'red':
        return RGBColor(0xff, 0x00, 0x00)
    else:
        print(f"WARNING: Color '{color}' not found in 'color_to_RGB'")
        return RGBColor(0x00, 0x00, 0x00)