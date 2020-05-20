import os
import glob
import json
import datetime
import pandas as pd
import numpy as np

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # pylint: disable=no-name-in-module

from libs.utils import fund_list_extractor, INDEXES
from libs.tools import trend_simple_forecast

from .slide_utils import slide_title_header, color_to_RGB, pptx_ui_errors
from .synopsis_slide import generate_synopsis_slide

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8

TEMP_DIR = os.path.join("output", "temp")


def make_fund_slides(prs, analysis: dict, **kwargs):
    """Make Fund Slides

    Arguments:
        prs {pptx obj} -- presentation
        analysis {dict} -- full data

    Optional Args:
        views {str} -- (default: {''})

    Returns:
        prs -- pptx presentation object
    """
    views = kwargs.get('views', '')
    funds = analysis.keys()
    for fund in funds:
        if fund != '_METRICS_':
            prs = add_fund_content(prs, fund, analysis, views=views)

    return prs


def add_fund_content(prs, fund: str, analysis: dict, **kwargs):
    """Add Fund Content

    Arguments:
        prs {pptx-object} -- powerpoint object
        fund {str} -- fund name
        analysis {dict} -- analysis dictionary of data content

    Optional Args:
        views {str} -- (default: {''})

    Returns:
        pptx-object -- modified pptx
    """
    views = kwargs.get('views', '')
    if views is None:
        return prs

    content_dir = os.path.join(TEMP_DIR, fund, views)

    # Title slide for a fund
    fund_name = INDEXES.get(fund, fund)

    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    top = Inches(0.1)
    left = Inches(4)
    width = Inches(5)
    height = Inches(2)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txtbox.text_frame

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f'{fund_name}'
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.name = 'Arial'

    p2 = text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = f"Dates Covered: {analysis[fund][views]['dates_covered']['start']}  :  " + \
        f"{analysis[fund][views]['dates_covered']['end']}"
    p2.font.bold = False
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x74, 0x3c, 0xe6)
    p2.font.name = 'Arial'

    has_beta = False

    # Insert a table of fund figures
    left_loc = Inches(0.1)
    top_loc = Inches(1.1)
    table_width = Inches(2.4)
    table_height = Inches(1.4)

    vq = analysis[fund].get('metadata', {}).get(
        'volatility', {}).get('VQ')
    has_vq = False
    rows = 8

    if vq is not None:
        has_vq = True
        rows = rows + 4
        stop_loss = analysis[fund].get('metadata', {}).get(
            'volatility', {}).get('stop_loss')
        high_close = analysis[fund].get('metadata', {}).get(
            'volatility', {}).get('last_max', {}).get('Price', 'n/a')
        status = analysis[fund].get('metadata', {}).get(
            'volatility', {}).get('status', {}).get('status', 'n/a')
        vq_color = analysis[fund].get('metadata', {}).get(
            'volatility', {}).get('status', {}).get('color', 'n/a')

    table_placeholder = slide.shapes.add_table(rows,
                                               2,
                                               left_loc,
                                               top_loc,
                                               table_width,
                                               table_height)
    table = table_placeholder.table

    table.cell(0, 0).text = 'Attribute'
    table.cell(0, 1).text = ''
    table.cell(1, 0).text = 'Current Price'
    table.cell(1, 1).text = '$' + \
        str(np.round(analysis[fund][views]
                     ['statistics']['current_price'], 2))

    table.cell(2, 0).text = 'Alpha'
    table.cell(3, 0).text = 'Beta'
    table.cell(4, 0).text = 'R-Squared'
    table.cell(5, 0).text = 'Sharpe Ratio'
    table.cell(6, 0).text = 'Std. Deviation'

    table.cell(2, 1).text = 'n/a'
    table.cell(3, 1).text = 'n/a'
    table.cell(4, 1).text = 'n/a'
    table.cell(5, 1).text = 'n/a'
    table.cell(6, 1).text = 'n/a'

    risk_ratios = analysis[fund][views]['statistics'].get('risk_ratios', {})
    if 'alpha' in risk_ratios:
        table.cell(2, 1).text = str(
            np.round(risk_ratios['alpha'], 5))
    if 'beta' in risk_ratios:
        table.cell(3, 1).text = str(
            np.round(risk_ratios['beta'], 5))
    if 'r_squared' in risk_ratios:
        table.cell(4, 1).text = str(
            np.round(risk_ratios['r_squared'], 5))
    if 'sharpe' in risk_ratios:
        table.cell(5, 1).text = str(
            np.round(risk_ratios['sharpe'], 5))
    if 'standard_deviation' in risk_ratios:
        table.cell(6, 1).text = str(
            np.round(risk_ratios['standard_deviation'], 5))

    table.cell(7, 0).text = 'Altman-Z Score'
    alt_z = analysis[fund].get('metadata', {}).get(
        'altman_z', {})
    alt_z_score = alt_z.get('score', "n/a")
    alt_z_color = alt_z.get('color', "black")

    if isinstance(alt_z_score, (float, int)):
        alt_z_score = str(np.round(alt_z_score, 5))

    table.cell(7, 1).text = alt_z_score
    table.cell(7, 1).text_frame.paragraphs[0].font.color.rgb = color_to_RGB(
        alt_z_color)

    end = 7

    if has_vq:
        table.cell(end+1, 0).text = 'Volatility Quotient'
        table.cell(end+1, 1).text = str(vq)
        table.cell(end+2, 0).text = 'Stop Loss'
        table.cell(end+2, 1).text = str(stop_loss)
        table.cell(end+3, 0).text = 'Last High Close'
        table.cell(end+3, 1).text = str(high_close)
        table.cell(end+4, 0).text = 'VQ Status'
        table.cell(end+4, 1).text = str(status)

        table.cell(end+4, 1).text_frame.paragraphs[0].font.color.rgb = color_to_RGB(
            vq_color)

    table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(16)
    table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(16)

    for i in range(1, rows):
        table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(14)
        table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(14)

    has_beta = True

    content = os.path.join(content_dir, f"candlestick_{fund}.png")
    if os.path.exists(content):
        if has_beta:
            left = Inches(2.6)  # Inches(1.42)
        else:
            left = Inches(1.42)

        top = Inches(1.4)
        height = Inches(6)
        width = Inches(10.5)
        slide.shapes.add_picture(
            content, left, top, height=height, width=width)

    else:
        slide = pptx_ui_errors(slide, "No Candlestick Chart available.")

    price_pt = np.round(analysis[fund][views]
                        ['statistics']['current_price'], 2)
    price_chg_p = np.round(
        analysis[fund][views]['statistics']['current_percent_change'], 3)
    price_chg = np.round(analysis[fund][views]
                         ['statistics']['current_change'], 2)

    if price_chg > 0.0:
        price_str = f"${price_pt} +{price_chg} (+{price_chg_p}%)"
    else:
        price_str = f"${price_pt} {price_chg} ({price_chg_p}%)"

    # Slide #0: synopsis
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    slide = generate_synopsis_slide(slide, analysis, fund, views=views)

    # Slide #1 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    indexes = []
    indexes.append(len(prs.slides) - 1)
    slide = slide_title_header(slide, fund, price_details=price_str)

    # Slide #2 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #3 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #4 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #5 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #6 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #7 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    # Slide #8 of content
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(slide, fund, price_details=price_str)
    indexes.append(len(prs.slides)-1)

    content = os.path.join(content_dir, '*.png')
    pics = glob.glob(content)

    content_dir2, _ = os.path.split(content_dir)
    content2 = os.path.join(content_dir2, '*.png')
    pics2 = glob.glob(content2)
    pics.extend(pics2)

    fund_analysis = analysis[fund]
    current_price = analysis[fund][views]['statistics']['current_price']
    prs = format_plots(prs, indexes, pics,
                       fund_analysis=fund_analysis, views=views, current_price=current_price)

    return prs


def format_plots(prs, slide_indices: list, globs: list, fund_analysis: dict = {}, **kwargs):
    """Format Plots

    Arguments:
        prs {pptx-object} -- entire presentation
        slide_indices {list} -- fund-specific slide numbers in a list for referencing
        globs {list} -- plot file path list

    Keyword Arguments:
        fund_analysis {dict} -- fund data object (default: {{}})

    Optional Args:
        views {str} -- (default: {''})
        current_price {float} -- (default: {None})

    Returns:
        pptx-object -- filled in slides with content
    """
    if len(globs) == 0:
        for ind in slide_indices:
            pptx_ui_errors(prs.slides[ind], "No plot files available.")
        return prs

    content_file = os.path.join(
        "libs", "ui_generation", "pptx_resources", "fund_content_slides.json")

    if not os.path.exists(content_file):
        for ind in slide_indices:
            pptx_ui_errors(
                prs.slides[ind], "File 'fund_content_slides.json' not found.")
        return prs

    slide_content = {}
    with open(content_file, 'r') as c_file:
        slide_content = json.load(c_file)
        c_file.close()

    views = kwargs.get('views', '')
    current_price = kwargs.get('current_price')

    locations = slide_content.get('locations', [])
    for picture in globs:

        _, part = os.path.split(picture)
        splits = part.split('_')
        splits.pop(-1)
        part = '_'.join(splits)

        if part in slide_content:
            details = slide_content.get(part, {})
            slide_index = details.get('index')
            location = details.get('location')

            left = eval(locations[location]['left'])
            top = eval(locations[location]['top'])
            height = eval(locations[location]['height'])
            width = eval(locations[location]['width'])

            prs.slides[slide_indices[slide_index]].shapes.add_picture(
                picture, left, top, height=height, width=width)

        # Slide #7
        slide_num = 6
        if 'resist_support' in part:
            left = Inches(0)
            top = Inches(1.55)
            height = Inches(4.7)
            width = Inches(7)
            prs.slides[slide_indices[slide_num]].shapes.add_picture(
                picture, left, top, height=height, width=width)

            left = Inches(7)
            top = Inches(0.25)
            height = Inches(4.7)
            width = Inches(4)
            txbox = prs.slides[slide_indices[slide_num]].shapes.add_textbox(
                left, top, width, height)

            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Nearest Support & Resistance Levels"
            p.font.size = Pt(18)
            p.font.name = 'Arial'
            p.font.bold = True

            left_loc = Inches(8)
            top_loc = Inches(1)
            table_width = Inches(4.5)

            num_srs = len(fund_analysis[views]
                          ['support_resistance']['major S&R']) + 1
            table_height = Inches(num_srs * 0.33)
            if num_srs * 0.33 > 6.0:
                table_height = Inches(6.0)

            table_placeholder = prs.slides[slide_indices[slide_num]].shapes.add_table(
                num_srs,
                3,
                left_loc,
                top_loc,
                table_width,
                table_height)
            table = table_placeholder.table

            table.cell(0, 0).text = 'Price'
            table.cell(0, 1).text = '% Change'
            table.cell(0, 2).text = 'Sprt / Res'

            for i, maj in enumerate(fund_analysis[views]['support_resistance']['major S&R']):
                table.cell(i+1, 0).text = f"${maj['Price']}"
                table.cell(i+1, 1).text = f"{maj['Change']}"
                table.cell(i+1, 2).text = maj['State']
                color = color_to_RGB(maj['Color'])
                table.cell(
                    i+1, 0).text_frame.paragraphs[0].font.color.rgb = color
                table.cell(
                    i+1, 1).text_frame.paragraphs[0].font.color.rgb = color
                table.cell(
                    i+1, 2).text_frame.paragraphs[0].font.color.rgb = color
                table.cell(
                    i+1, 0).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(
                    i+1, 1).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(
                    i+1, 2).text_frame.paragraphs[0].font.size = Pt(14)

        # Slide 8
        slide_num = 7
        if 'trendline' in part:
            left = Inches(0)
            top = Inches(1.55)
            height = Inches(4.7)
            width = Inches(7)
            prs.slides[slide_indices[slide_num]].shapes.add_picture(
                picture, left, top, height=height, width=width)

            left = Inches(6.5)
            top = Inches(0.25)
            height = Inches(4.7)
            width = Inches(4)
            txbox = prs.slides[slide_indices[slide_num]].shapes.add_textbox(
                left, top, width, height)

            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Long, Intermediate, Short, and Near Term Trendlines"
            p.font.size = Pt(18)
            p.font.name = 'Arial'
            p.font.bold = True

            trends = []
            futures = [0, 10, 20, 40, 60, 80]
            forecasts = []
            for trend in fund_analysis[views]['trendlines']['current']:
                trends.append(trend)
                forecast = trend_simple_forecast(
                    trend, future_periods=futures, current_price=current_price)
                forecasts.append(forecast)

            num_rows = len(trends) + 2
            num_cols = len(futures) + 2
            table_height = num_rows * 0.33
            table_width = num_cols * 0.8
            if num_rows * 0.33 > 6.0:
                table_height = 6.0
            if num_cols * 0.8 > 6.6:
                table_width = 6.6

            left_start = 6.5
            adj_left = (6.6 - table_width) / 2
            table_height = Inches(table_height)
            table_width = Inches(table_width)
            left_loc = Inches(left_start + adj_left)
            top_loc = Inches(0.75)

            table_placeholder = prs.slides[slide_indices[slide_num]].shapes.add_table(
                num_rows,
                num_cols,
                left_loc,
                top_loc,
                table_width,
                table_height)
            table = table_placeholder.table

            cell_1 = table.cell(0, 0)
            cell_2 = table.cell(0, num_cols-1)
            cell_1.merge(cell_2)
            cell_1.text = f"Future Periods of Active Trendlines"

            for i, fut in enumerate(futures):
                table.cell(1, i+1).text = str(fut)
                table.cell(1, i+1).text_frame.paragraphs[0].font.bold = True

            table.cell(1, 0).text = 'Trend'
            table.cell(1, 0).text_frame.paragraphs[0].font.bold = True

            table.cell(1, i+2).text = 'Price'
            table.cell(1, i+2).text_frame.paragraphs[0].font.bold = True

            for i, trend in enumerate(trends):
                for j, value in enumerate(forecasts[i]['returns']):
                    table.cell(i+2, j+1).text = f"${value}"
                    table.cell(
                        i+2, j+1).text_frame.paragraphs[0].font.size = Pt(12)
                    color = color_to_RGB(trend['color'])
                    table.cell(
                        i+2, j+1).text_frame.paragraphs[0].font.color.rgb = color

                table.cell(i+2, 0).text = forecasts[i]['slope']
                table.cell(
                    i+2, 0).text_frame.paragraphs[0].font.size = Pt(12)
                color = color_to_RGB(trend['color'])
                table.cell(
                    i+2, 0).text_frame.paragraphs[0].font.color.rgb = color

                table.cell(i+2, j+2).text = forecasts[i]['above_below']
                table.cell(
                    i+2, j+2).text_frame.paragraphs[0].font.size = Pt(12)
                color = color_to_RGB(trend['color'])
                table.cell(
                    i+2, j+2).text_frame.paragraphs[0].font.color.rgb = color

    return prs
