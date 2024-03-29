""" Synopsis Slide Generator """
import numpy as np

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN  # pylint: disable=no-name-in-module

from .slide_utils import pptx_ui_errors, COLOR_TO_RGB, space_injector

ORDER_RANK = {
    "hull_moving_average": 1,
    "exponential_moving_average": 10,
    "simple_moving_average": 100
}


def generate_synopsis_slide(slide, analysis: dict, fund: str, **kwargs):
    """Generate Synopsis Slide

    Arguments:
        slide {pptx-slide} -- slide object
        analysis {dict} -- entire data object
        fund {str} -- fund name

    Optional Args:
        views {str} -- (default: {None})

    Returns:
        pptx-slide -- modified slide object
    """
    views = kwargs.get('views')
    if views is None:
        return pptx_ui_errors(slide, "No 'views' object passed.")

    synopsis = analysis[fund]['synopsis'][views]
    slide = add_synopsis_title(slide, f"Metrics Summary ({views})")

    slide = add_synopsis_category_box(
        slide, 'trend', synopsis, type_='metrics')
    slide = add_synopsis_category_box(
        slide, 'oscillator', synopsis, type_='metrics')
    slide = add_synopsis_category_box(
        slide, 'trendlines', synopsis, type_='metrics')

    return slide


def add_synopsis_title(slide, title: str):
    """Synopsis Title

    Add the title to the synopsis slide

    Arguments:
        slide {pptx-slide} -- slide object
        title {str} -- title to add

    Returns:
        pptx-slide -- modified slide object
    """
    top = Inches(0)
    left = Inches(7.55)
    width = Inches(5)
    height = Inches(2)
    txt_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txt_box.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = f'{title}'
    paragraph.font.bold = True
    paragraph.font.size = Pt(40)
    paragraph.font.name = 'Arial'

    return slide


def add_synopsis_category_box(slide, category: str, content: dict, type_='metrics'):
    """Add Synopsis Category Box

    Arguments:
        slide {pptx-slide} -- slide object
        category {str} -- category name to add
        content {dict} -- content to add

    Keyword Arguments:
        type_ {str} -- controlling key name (default: {'metrics'})

    Returns:
        pptx-slide -- modified slide object
    """
    # pylint: disable=too-many-locals,too-many-branches,too-many-statements
    cat = type_ + "_categories"
    listed = content.get(cat, {}).get(category, [])

    top_title_inches = 0.9
    top_table_inches = 1.35
    line_length = 27

    if len(listed) == 0:
        # Trend line case... look in metrics alone:
        listed = content.get(type_, {}).get(category, [])

    if type_ == 'metrics':
        if category == 'trend':
            listed = reorder_trends(listed)

            # Trend Metrics Table
            top = Inches(top_title_inches)
            left = Inches(0.15)
            width = Inches(4.25)
            height = Inches(.58)
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = txt_box.text_frame

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.text = 'Trend-Based Metrics'
            paragraph.font.bold = True
            paragraph.font.size = Pt(20)
            paragraph.font.name = 'Arial'

            table_height = Inches(3)
            table_width = Inches(4.25)
            left_loc = left
            top_loc = Inches(top_table_inches)
            table_placeholder = slide.shapes.add_table(len(listed) + 1,
                                                       2,
                                                       left_loc,
                                                       top_loc,
                                                       table_width,
                                                       table_height)
            table = table_placeholder.table

            table.cell(0, 0).text = "Current vs. Metric"
            table.cell(0, 1).text = "Current\t(Previous)"
            table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(14)
            table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(14)

            for i, item in enumerate(listed):
                item_str = pretty_up_key(item)
                value = content.get('metrics', {}).get(item, 0.0)
                delta = content.get('metrics_delta', {}).get(item, 0.0)

                value_str = f"{value}%"
                delta_str = f"({delta}%)"

                if value > 0.0:
                    value_str = f"+{value}%"
                if delta > 0.0:
                    delta_str = f"(+{delta}%)"

                value_str = value_str + ' ' + delta_str
                value_str = space_injector(value_str, line_length)

                table.cell(i+1, 0).text = item_str
                table.cell(i+1, 1).text = value_str
                table.cell(i+1, 0).text_frame.paragraphs[0].font.size = Pt(10)
                table.cell(i+1, 1).text_frame.paragraphs[0].font.size = Pt(10)

                if value > 0.0:
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["green"]
                elif value < 0.0:
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["red"]

            # A note about moving averages
            note = "'Current vs. Metric' refers to difference between moving average metric " + \
                "and current close. A negative % refers to a close X% less than the metric. " + \
                "(Previous) is the previous period's close percent difference. This is " + \
                "supplied to see change to current day."

            top = Inches(6.8)
            width = Inches(4.25)
            height = Inches(0.56)
            txt_box2 = slide.shapes.add_textbox(left, top, width, height)
            text_frame = txt_box2.text_frame
            text_frame.word_wrap = True

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.text = note
            paragraph.font.bold = False
            paragraph.font.size = Pt(8)
            paragraph.font.name = 'Arial'

        elif category == 'oscillator':
            # Oscillator Metrics Table
            top = Inches(top_title_inches)
            left = Inches(4.55)
            width = Inches(4.25)
            height = Inches(.58)
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = txt_box.text_frame

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.text = 'Oscillator-Based Metrics'
            paragraph.font.bold = True
            paragraph.font.size = Pt(20)
            paragraph.font.name = 'Arial'

            table_height = Inches(4)
            table_width = Inches(4.25)
            left_loc = left
            top_loc = Inches(top_table_inches)
            table_placeholder = slide.shapes.add_table(len(listed) + 1,
                                                       2,
                                                       left_loc,
                                                       top_loc,
                                                       table_width,
                                                       table_height)
            table = table_placeholder.table

            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Current\t(Previous)"
            table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(14)
            table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(14)

            for i, item in enumerate(listed):
                item_str = pretty_up_key(item)
                value = content.get('metrics', {}).get(item, 0.0)
                delta = content.get('metrics_delta', {}).get(item, 0.0)

                value_str = f"{np.round(value, 5)}"
                delta_str = f"({np.round(delta, 5)})"
                value_str = value_str + ' ' + delta_str
                value_str = space_injector(value_str, line_length)

                table.cell(i+1, 0).text = item_str
                table.cell(i+1, 1).text = value_str
                table.cell(i+1, 0).text_frame.paragraphs[0].font.size = Pt(10)
                table.cell(i+1, 1).text_frame.paragraphs[0].font.size = Pt(10)

                if value > 0.0:
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["green"]
                elif value < 0.0:
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["red"]

        elif category == 'trendlines':
            # Trendlines Metrics Table
            top = Inches(top_title_inches)
            left = Inches(8.95)
            width = Inches(4.25)
            height = Inches(.58)
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = txt_box.text_frame

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.text = 'Current Calculated Trend Lines'
            paragraph.font.bold = True
            paragraph.font.size = Pt(20)
            paragraph.font.name = 'Arial'

            height = min((float(len(listed)) + 1.0) * 0.4, 4.0)
            table_height = Inches(height)
            table_width = Inches(4.25)
            left_loc = left
            top_loc = Inches(top_table_inches)
            table_placeholder = slide.shapes.add_table(len(listed) + 1,
                                                       2,
                                                       left_loc,
                                                       top_loc,
                                                       table_width,
                                                       table_height)
            table = table_placeholder.table

            table.cell(0, 0).text = "Trend (Length)"
            table.cell(0, 1).text = "Type"
            table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(14)
            table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(14)

            for i, item in enumerate(listed):
                periods = item.get('periods', 0)
                term = ''
                for key in item:
                    if key != 'periods':
                        term = key

                style = item.get(term, '').capitalize()
                term_str = pretty_up_key(term, parser=' ')
                term_str = f"{term_str} ({periods})"

                table.cell(i+1, 0).text = term_str
                table.cell(i+1, 1).text = style
                table.cell(i+1, 0).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+1, 1).text_frame.paragraphs[0].font.size = Pt(14)
                table.cell(i+1, 1).text_frame.paragraphs[0].font.bold = True

                if style == 'Bull':
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["green"]
                else:
                    table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = \
                        COLOR_TO_RGB["red"]

    return slide


def pretty_up_key(key: str, parser: str = '_') -> str:
    """Pretty-Up Key

    Arguments:
        key {str} -- key to "make pretty"

    Keyword Arguments:
        parser {str} -- char to split key (default: {'_'})

    Returns:
        str -- new key, capitalized and w/o "parser" in it
    """
    keys = key.split(parser)
    keys = [new_key.capitalize() for new_key in keys]
    output = ' '.join(keys)
    return output


def reorder_trends(listed: list) -> list:
    """Reorder Trends

    Argumentss:
        listed {list} -- list of trends to reorder based on time

    Returns:
        list -- reordered list
    """
    new_list = []
    content = []
    non_timing = []
    for item in listed:
        split = item.split('(')

        if '-d' in split[1]:
            period = split[1].split('-')[0]
            data = {"period": int(period), "type": split[0]}
            content.append(data)
        else:
            non_timing.append(item)

    content.sort(key=lambda x: x['period'])
    for item in content:
        val = '('.join([item['type'], str(item['period'])])
        val += '-d)'
        new_list.append(val)

    for item in non_timing:
        new_list.append(item)

    return new_list
