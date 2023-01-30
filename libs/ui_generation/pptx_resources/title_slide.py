""" title slide creator """
from datetime import datetime

from pptx.util import Inches, Pt
from pptx.presentation import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # pylint: disable=no-name-in-module

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8


def create_presentation_title(prs: Presentation, version: str) -> Presentation:
    """Title Slide

    Arguments:
        prs {pptx-obj} -- powerpoint python object
        VERSION {str} -- '0.2.02', for example

    Returns:
        pptx-obj -- powerpoint object
    """
    height = prs.slide_height
    width = int(16 * height / 9)
    prs.slide_width = width
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])

    left_inches = 6
    left = Inches(left_inches)
    top = Inches(2.45)
    text = slide.shapes.add_textbox(left, top, Inches(1), Inches(1))
    text_frame = text.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = 'Securities Analysis'
    paragraph.font.bold = True
    paragraph.font.size = Pt(48)
    paragraph.font.name = 'Arial'

    paragraph4 = text_frame.add_paragraph()
    paragraph4.alignment = PP_ALIGN.CENTER
    paragraph4.text = "A Technical Analysis of Financial Markets by 'nga-27'"
    paragraph4.font.italic = True
    paragraph4.font.size = Pt(14)
    paragraph4.font.color.rgb = RGBColor(0x74, 0x3c, 0xe6)
    paragraph4.font.name = 'Arial'

    left = Inches(left_inches)
    top = Inches(4.0)
    text = slide.shapes.add_textbox(left, top, Inches(1), Inches(1))
    text_frame2 = text.text_frame

    paragraph2 = text_frame2.paragraphs[0]
    paragraph2.alignment = PP_ALIGN.CENTER
    paragraph2.text = f'Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}'
    paragraph2.font.bold = False
    paragraph2.font.size = Pt(22)
    paragraph2.font.color.rgb = RGBColor(0x30, 0x9c, 0x4f)
    paragraph2.font.name = 'Arial'

    paragraph3 = text_frame2.add_paragraph()
    paragraph3.alignment = PP_ALIGN.CENTER
    paragraph3.text = f'Software Version: {version}'
    paragraph3.font.bold = False
    paragraph3.font.size = Pt(18)
    paragraph3.font.color.rgb = RGBColor(0x30, 0x9c, 0x4f)
    paragraph3.font.name = 'Arial'

    return prs
