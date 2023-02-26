""" slide utilities """
import os
from datetime import datetime

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # pylint: disable=no-name-in-module

from libs.utils import INDEXES, STANDARD_COLORS

# Slide Layouts
PRES_TITLE_SLIDE = 0
TITLE_CONTENT_SLIDE = 1
SECTION_HEADER_SLIDE = 2
TWO_CONTENT_SLIDE = 3
COMPARISON_SLIDE = 4
TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
CONTENT_W_CAPTION_SLIDE = 7
PICTURE_W_CAPTION_SLIDE = 8

WARNING = STANDARD_COLORS["warning"]
NORMAL = STANDARD_COLORS["normal"]

COLOR_TO_RGB = {
    'black': RGBColor(0x00, 0x00, 0x00),
    'blue': RGBColor(0x00, 0x00, 0xFF),
    'green': RGBColor(0x00, 0xCC, 0x00),
    'purple': RGBColor(0x7F, 0x00, 0xFF),
    'yellow':RGBColor(0xee, 0xd4, 0x00),
    'orange': RGBColor(0xff, 0x99, 0x33),
    'red': RGBColor(0xff, 0x00, 0x00)
}


def slide_title_header(slide, fund: str, include_time=True, price_details=''):
    """Slide Title Header

    Generates a consistent title header for fund slides.

    Arguments:
        slide {pptx-slide} -- slide object
        fund {str} -- fund name

    Keyword Arguments:
        include_time {bool} -- will include time of pptx creation (default: {True})
        price_details {str} -- price content for latest period (default: {''})

    Returns:
        pptx-slide -- updated slide with title content
    """
    fund = INDEXES.get(fund, fund)

    left = Inches(0)  # Inches(3.86)
    top = Inches(0)
    height = Inches(0.5)
    width = Inches(0.5)
    txt_box = slide.shapes.add_textbox(left, top, width, height)
    txt_frame = txt_box.text_frame

    paragraph = txt_frame.paragraphs[0]
    paragraph.text = fund
    paragraph.font.size = Pt(36)
    paragraph.font.name = 'Arial'
    paragraph.font.bold = True

    if include_time:
        paragraph = txt_frame.add_paragraph()
        paragraph.font.size = Pt(14)
        paragraph.font.bold = False
        paragraph.font.name = 'Arial'
        paragraph.text = str(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    if price_details != '':
        left = Inches(2)
        top = Inches(0.1)
        height = Inches(0.5)
        width = Inches(0.5)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        txt_frame = txt_box.text_frame

        paragraph = txt_frame.paragraphs[0]
        paragraph.text = fund
        paragraph.font.size = Pt(24)
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True
        paragraph.text = price_details
        details = price_details.split(' ')

        if '+' in details[1]:
            paragraph.font.color.rgb = COLOR_TO_RGB['green']
        else:
            paragraph.font.color.rgb = COLOR_TO_RGB['red']

    return slide


def subtitle_header(slide, title: str):
    """Subtitle Header

    Creates subtitle under main slide title

    Arguments:
        slide {pptx-slide} -- slide object
        title {str} -- title to add as subtitle

    Returns:
        pptx-slide -- modified slide object
    """
    top = Inches(0.61)
    left = Inches(0.42)
    height = Inches(0.5)
    width = Inches(0.5)
    txt_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txt_box.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.bold = False
    paragraph.font.size = Pt(22)
    paragraph.font.name = 'Times New Roman'

    return slide


def intro_slide(prs):
    """Intro Slide

    Adds some introductory slides to presentation, regarding metrics, etc.

    Arguments:
        prs {pptx-object} -- presentation

    Returns:
        pptx-object -- modified presentation
    """
    # pylint: disable=too-many-locals
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    slide = slide_title_header(
        slide, 'Explanation of Analysis', include_time=False)

    explanation_file_path = os.path.join("resources", "metric_explanation.txt")

    if os.path.exists(explanation_file_path):
        with open(explanation_file_path, 'r', encoding='utf-8') as filer:
            content = filer.readlines()

        content2 = []
        for cont in content:
            cropped = cont.split('\r\n')[0]
            content2.append(cropped)
        content = content2
        filer.close()

        top = Inches(0.81)
        left = Inches(0.42)
        width = Inches(11)
        height = Inches(6)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txt_box.text_frame
        text_frame.word_wrap = True

        paragraph = text_frame.paragraphs[0]
        paragraph.text = content[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True

        for i in range(1, len(content)):
            paragraph = text_frame.add_paragraph()
            paragraph.text = content[i]

            if i == 5:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
            else:
                paragraph.font.size = Pt(10)
                paragraph.font.bold = False

    else:
        print(f"{WARNING}WARNING - file 'metric_explanation.txt' not found.{NORMAL}")

    return prs


def pptx_ui_errors(slide, message: str):
    """PPTX UI Errors

    Log errors as they occur, print on slide itself

    Arguments:
        slide {pptx-slide} -- pptx-slide object, for logging purposes
        message {str} -- message to log

    Returns:
        pptx-slide -- slide with error listed on it
    """
    top = Inches(2.75)
    left = Inches(4)
    width = Inches(6)
    height = Inches(2)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txtbox.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = f'{message}'
    paragraph.font.bold = True
    paragraph.font.size = Pt(40)
    paragraph.font.name = 'Arial'

    return slide


def space_injector(space_sep_str: str, desired_str_len: int, sep=' ') -> str:
    """Space Injector

    Inject spaces or 'sep' between space-separated strings so that the end string length
    matches a specific size (used for printing and UI).

    Arguments:
        space_sep_str {str} -- string that is space separated of which to fix to a length
        desired_str_len {int} -- desired string length after adjustment

    Keyword Arguments
        sep {str} -- separator between strings (default: {' '})

    Returns:
        str -- newly spaced string
    """
    sep_str = space_sep_str.split(sep)
    sum_len = 0
    for sstr in sep_str:
        sum_len += len(sstr)

    diff = desired_str_len - sum_len
    diff_adj = int(float(diff) / float((len(sep_str) - 1)))
    new_sep = ' ' * diff_adj
    new_str = new_sep.join(sep_str)

    return new_str


class LocationType:
    """ Pass the location data from fund_content_slides.json and generate an object """
    # pylint: disable=too-few-public-methods
    left: Inches
    top: Inches
    height: Inches
    width: Inches

    def __init__(self):
        self.left = 0.0
        self.top = 0.0
        self.height = 0.0
        self.width = 0.0


def get_locations(location_dict: dict) -> LocationType:
    """get_locations

    Args:
        location_dict (dict): from fund_content_slides.json

    Returns:
        LocationType: Location data
    """
    location = LocationType()
    location.left = Inches(location_dict['left'])
    location.top = Inches(location_dict['top'])
    location.height = Inches(location_dict['height'])
    location.width = Inches(location_dict['width'])
    return location
